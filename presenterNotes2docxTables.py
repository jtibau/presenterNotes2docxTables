import sys
import collections 
import collections.abc
import pptx
from pptx import Presentation
from docx import Document
from docx.shared import Inches

# We process the pptx first, the code is more readable

# if(sys.argv[1]):
try:
    filename = sys.argv[1]
    ppt=Presentation(filename)
except IndexError:
    print("Provide the filename as a parameter")
    sys.exit()
except pptx.exc.PackageNotFoundError:
    print(filename+": File not Found or not a PPTX")
    sys.exit()



notes = []
for page, slide in enumerate(ppt.slides):
    textNote = slide.notes_slide.notes_text_frame.text
    # The first slide has some metadata info
    # that we'll throw in the header of the word doc
    if (page==0):
        headerText = textNote
        print(headerText)
    else:
        notes.append((page,textNote))


# The pptx is done, we'll start creating the word doc
document = Document()

# The header is what was present on the first slide
header = document.sections[0].header
header.paragraphs[0].text  = headerText

# Setting up the table
table = document.add_table(rows=1, cols=3)
table.style = "Table Grid"
table.allow_autofit = True
table.columns[0].width = Inches(1)
table.columns[1].width = Inches(5)
table.columns[2].width = Inches(1)

# Setting the Table Heading
hdr_cells = table.rows[0].cells
hdr_cells[0].text = 'OBJETOS DE VIDEO'
hdr_cells[1].text = 'ANOTACIONES'
hdr_cells[2].text = 'TOMA'

# Throw in all the presenter's notes into rows
for slideNumber, script in notes:
    if script:
        row_cells = table.add_row().cells
        row_cells[0].text = str(slideNumber+1)
        row_cells[1].text = script.strip()
        row_cells[2].text = ""

outputFilename = filename[:-5]+".docx"
document.save(outputFilename)